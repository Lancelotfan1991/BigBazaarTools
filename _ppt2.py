# -*- coding: utf-8 -*-
"""大巴扎趣闻录 PPT v2 — DB-driven: correct hero attribution + item images + effects."""
import os, io, json, hashlib, urllib.request
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

BASE = os.path.dirname(os.path.abspath(__file__))
IMGDIR = os.path.join(BASE, 'output', '_imgs')
os.makedirs(IMGDIR, exist_ok=True)

# ---- DB index ----
_all = json.load(open(os.path.join(BASE, 'mobile-app/data-s16.1/All.json'), encoding='utf-8'))
ITEMS = {it['名称']: it for it in _all['物品']}
SKILLS = {sk['名称']: sk for sk in _all.get('技能', [])}
SKILLS_EN = {(sk.get('英文名') or '').lower(): sk for sk in _all.get('技能', [])}
ITEMS_EN = {(it.get('英文名') or '').lower(): it for it in _all['物品']}

# ---- Theme ----
NAVY = RGBColor(0x0B, 0x1D, 0x3A)
PANEL = RGBColor(0x12, 0x2A, 0x4D)
CARDBG = RGBColor(0x0E, 0x22, 0x42)
GOLD = RGBColor(0xE8, 0xB8, 0x4B)
SEA = RGBColor(0x6F, 0xB3, 0xD9)
WHITE = RGBColor(0xF5, 0xF7, 0xFA)
MUTED = RGBColor(0xAE, 0xBE, 0xD0)
FONT = "Microsoft YaHei"

# hero -> (display, color)
HERO = {
    '海盗': ('海盗 · 瓦妮莎', RGBColor(0x2E, 0x9C, 0xC4)),
    '杜利': ('杜利', RGBColor(0xE0, 0x7B, 0x39)),
    '朱尔斯': ('朱尔斯', RGBColor(0xD9, 0x5B, 0x8F)),
    '皮格马利安': ('皮格马利安', RGBColor(0x4F, 0xA8, 0x5E)),
    '斯特尔': ('斯特尔', RGBColor(0x8B, 0x6F, 0xC9)),
    '马克': ('马克', RGBColor(0x9B, 0x59, 0xB6)),
    '卡诺克': ('卡诺克', RGBColor(0xC0, 0x4A, 0x3F)),
    '通用': ('通用', RGBColor(0x7A, 0x8A, 0x9A)),
}

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def _set_font(run, name=FONT):
    run.font.name = name
    rPr = run.font._rPr
    if rPr is not None:
        ea = rPr.find(qn('a:ea'))
        if ea is None:
            ea = rPr.makeelement(qn('a:ea'), {}); rPr.append(ea)
        ea.set('typeface', name)


def add_bg(slide, color=NAVY):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, x, y, w, h, color, line=None, rounded=False):
    shp = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sp = slide.shapes.add_shape(shp, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.06); tf.margin_right = Inches(0.06)
    tf.margin_top = 0; tf.margin_bottom = 0
    return tf


def add_para(tf, text, size, color, bold=False, first=False,
             align=PP_ALIGN.LEFT, space_after=8):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align; p.space_after = Pt(space_after); p.space_before = Pt(0)
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    _set_font(r)
    return p


def add_runs(tf, runs, size, first=False, space_after=8, align=PP_ALIGN.LEFT):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align; p.space_after = Pt(space_after); p.space_before = Pt(0)
    for text, color, bold in runs:
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
        _set_font(r)
    return p


# ---- DB lookup + image cache ----
def lookup(name):
    """Return (entry, kind) where kind in item/skill; None if missing."""
    if name in ITEMS:
        return ITEMS[name], 'item'
    if name in SKILLS:
        return SKILLS[name], 'skill'
    low = name.lower()
    if low in ITEMS_EN:
        return ITEMS_EN[low], 'item'
    if low in SKILLS_EN:
        return SKILLS_EN[low], 'skill'
    return None, None


def owner_of(entry):
    ow = entry.get('所属职业') or []
    return ow[0] if ow else '通用'


def effect_lines(entry):
    out = []
    for e in entry.get('效果说明', []) or []:
        t = (e.get('text') or '').strip()
        if not t:
            continue
        cond = (e.get('condition') or '').strip()
        out.append(t + (f"（{cond}）" if cond else ''))
    return out


def fetch_img(entry):
    url = (entry or {}).get('图片链接') or ''
    if not url:
        return None
    h = hashlib.sha1(url.encode()).hexdigest()[:16]
    path = os.path.join(IMGDIR, h + '.png')
    if os.path.exists(path):
        return path
    try:
        req = urllib.request.Request(url, headers={
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) '
                          'AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120 Safari/537.36',
            'Referer': 'https://bazaardb.gg/'})
        data = urllib.request.urlopen(req, timeout=25).read()
        im = Image.open(io.BytesIO(data)).convert('RGBA')
        im.save(path)
        return path
    except Exception as e:
        print('  [img fail]', url, type(e).__name__, e)
        return None


def add_image_fit(slide, path, box_x, box_y, box_w, box_h):
    """Place image contained within box, centered, preserving aspect."""
    with Image.open(path) as im:
        iw, ih = im.size
    bw, bh = float(box_w), float(box_h)
    scale = min(bw / iw, bh / ih)
    w = int(iw * scale); h = int(ih * scale)
    x = int(box_x + (bw - w) / 2); y = int(box_y + (bh - h) / 2)
    slide.shapes.add_picture(path, Emu(x), Emu(y), Emu(w), Emu(h))


def _kicker(slide, text):
    if not text:
        return
    tf = textbox(slide, Inches(0.7), Inches(0.4), Inches(12), Inches(0.4))
    add_para(tf, text, 14, SEA, bold=True, first=True)


def _badge(slide, x, y, text, color):
    w = Inches(0.5 + 0.16 * len(text))
    sp = rect(slide, x, y, w, Inches(0.42), color, rounded=True)
    tf = sp.text_frame
    tf.word_wrap = False
    tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = WHITE
    _set_font(r)
    return w


def title_slide(title, subtitle, tagline):
    s = prs.slides.add_slide(BLANK); add_bg(s)
    rect(s, Inches(0), Inches(2.55), Inches(13.333), Inches(0.06), GOLD)
    rect(s, Inches(0), Inches(4.35), Inches(13.333), Inches(0.06), GOLD)
    tf = textbox(s, Inches(1), Inches(2.7), Inches(11.333), Inches(1.6), MSO_ANCHOR.MIDDLE)
    add_para(tf, title, 54, GOLD, bold=True, first=True, align=PP_ALIGN.CENTER, space_after=6)
    tf2 = textbox(s, Inches(1), Inches(4.5), Inches(11.333), Inches(1.2))
    add_para(tf2, subtitle, 22, WHITE, first=True, align=PP_ALIGN.CENTER, space_after=10)
    add_para(tf2, tagline, 15, MUTED, align=PP_ALIGN.CENTER)
    return s


def section_slide(num, title, subtitle):
    s = prs.slides.add_slide(BLANK); add_bg(s, PANEL)
    rect(s, Inches(1.0), Inches(2.75), Inches(0.18), Inches(2.0), GOLD)
    tf = textbox(s, Inches(1.45), Inches(2.7), Inches(10.5), Inches(2.2), MSO_ANCHOR.MIDDLE)
    add_para(tf, num, 20, SEA, bold=True, first=True, space_after=6)
    add_para(tf, title, 40, WHITE, bold=True, space_after=8)
    if subtitle:
        add_para(tf, subtitle, 18, MUTED)
    return s


def item_slide(name, kicker, anecdote, en_override=None, effect_override=None):
    """DB-driven: image (left) + hero badge/name/effect/anecdote (right)."""
    entry, kind = lookup(name)
    s = prs.slides.add_slide(BLANK); add_bg(s)
    _kicker(s, kicker)
    # left card
    cx, cy, cw, ch = Inches(0.7), Inches(1.25), Inches(3.5), Inches(5.6)
    rect(s, cx, cy, cw, ch, CARDBG, rounded=True)
    img = fetch_img(entry) if entry else None
    if img:
        add_image_fit(s, img, int(cx) + Inches(0.18), int(cy) + Inches(0.18),
                      int(cw) - Inches(0.36), int(ch) - Inches(0.36))
    else:
        tf = textbox(s, cx, cy, cw, ch, MSO_ANCHOR.MIDDLE)
        add_para(tf, '（暂无卡面）', 15, MUTED, first=True, align=PP_ALIGN.CENTER)
    # right column
    rx = Inches(4.55); rw = Inches(8.0)
    owner = owner_of(entry) if entry else '通用'
    disp, col = HERO.get(owner, (owner, MUTED))
    _badge(s, rx, Inches(1.25), disp, col)
    en = en_override or (entry.get('英文名') if entry else '') or ''
    tf = textbox(s, rx, Inches(1.85), rw, Inches(1.15))
    add_para(tf, name, 34, GOLD, bold=True, first=True, space_after=2)
    if en:
        add_para(tf, en, 18, MUTED)
    rect(s, rx, Inches(3.05), rw, Inches(0.02), SEA)
    # effect
    effs = effect_override if effect_override is not None else (effect_lines(entry) if entry else [])
    tf = textbox(s, rx, Inches(3.2), rw, Inches(1.85))
    add_para(tf, '效果说明', 16, SEA, bold=True, first=True, space_after=6)
    if effs:
        for e in effs[:5]:
            add_runs(tf, [('•  ', SEA, True), (e, WHITE, False)], 14, space_after=6)
    else:
        add_para(tf, '（该物品在当前数据版本中无效果文本）', 13, MUTED)
    # anecdote
    tf = textbox(s, rx, Inches(5.15), rw, Inches(1.95))
    add_para(tf, '趣闻', 16, GOLD, bold=True, first=True, space_after=6)
    for a in anecdote:
        add_para(tf, a, 15, WHITE, space_after=7)
    return s


def grid_slide(kicker, title, subtitle, names, notes):
    """Grid of small item cards: image + zh/en + one-line note. notes: dict."""
    s = prs.slides.add_slide(BLANK); add_bg(s)
    _kicker(s, kicker)
    tf = textbox(s, Inches(0.7), Inches(0.75), Inches(12), Inches(0.95))
    add_para(tf, title, 30, GOLD, bold=True, first=True, space_after=2)
    if subtitle:
        add_para(tf, subtitle, 15, MUTED)
    rect(s, Inches(0.7), Inches(1.78), Inches(11.9), Inches(0.02), SEA)
    n = len(names)
    cols = min(n, 4)
    rows = (n + cols - 1) // cols
    top0 = 2.0
    cellw = 11.9 / cols
    cellh = (7.2 - top0) / rows
    for i, nm in enumerate(names):
        entry, _ = lookup(nm)
        r = i // cols; c = i % cols
        x = 0.7 + c * cellw; y = top0 + r * cellh
        imgh = min(cellh - 1.35, 2.1)
        img = fetch_img(entry) if entry else None
        if img:
            add_image_fit(s, img, Inches(x), Inches(y), Inches(cellw), Inches(imgh))
        owner = owner_of(entry) if entry else '通用'
        _, col = HERO.get(owner, (owner, MUTED))
        tf = textbox(s, Inches(x + 0.08), Inches(y + imgh + 0.02), Inches(cellw - 0.16), Inches(cellh - imgh - 0.04))
        en = (entry.get('英文名') if entry else '') or ''
        add_runs(tf, [(nm, GOLD, True)], 15, first=True, space_after=1, align=PP_ALIGN.CENTER)
        if en:
            add_para(tf, en, 11, col, align=PP_ALIGN.CENTER, space_after=2)
        note = notes.get(nm, '')
        if note:
            add_para(tf, note, 12, MUTED, align=PP_ALIGN.CENTER, space_after=0)
    return s


def note_slide(kicker, title, en, blocks, table=None):
    """Text slide for translation topics (no single item). blocks: list of
       ('h'|'b'|'p'|'q'|'runs', val). table: optional (headers, rows)."""
    s = prs.slides.add_slide(BLANK); add_bg(s)
    rect(s, Inches(0.7), Inches(0.7), Inches(0.16), Inches(1.15), GOLD)
    tf = textbox(s, Inches(1.05), Inches(0.6), Inches(11.5), Inches(1.4))
    if kicker:
        add_para(tf, kicker, 14, SEA, bold=True, first=True, space_after=4)
        add_runs(tf, [(title, GOLD, True), ('   ' + en if en else '', MUTED, False)], 30)
    else:
        add_runs(tf, [(title, GOLD, True), ('   ' + en if en else '', MUTED, False)], 30, first=True)
    rect(s, Inches(1.05), Inches(1.95), Inches(11.5), Inches(0.025), SEA)
    tf = textbox(s, Inches(1.05), Inches(2.2), Inches(11.5), Inches(4.9))
    first = True
    for kind, val in blocks:
        if kind == 'h':
            add_para(tf, val, 19, GOLD, bold=True, first=first, space_after=6)
        elif kind == 'b':
            add_runs(tf, [('•  ', SEA, True), (val, WHITE, False)], 16, first=first, space_after=8)
        elif kind == 'p':
            add_para(tf, val, 17, WHITE, first=first, space_after=9)
        elif kind == 'q':
            add_para(tf, '“' + val + '”', 16, SEA, first=first, space_after=9)
        elif kind == 'runs':
            add_runs(tf, val, 16, first=first, space_after=8)
        first = False
    return s


def closing_slide(title, lines):
    s = prs.slides.add_slide(BLANK); add_bg(s, PANEL)
    rect(s, Inches(0), Inches(3.0), Inches(13.333), Inches(0.06), GOLD)
    tf = textbox(s, Inches(1), Inches(2.4), Inches(11.333), Inches(0.9), MSO_ANCHOR.MIDDLE)
    add_para(tf, title, 40, GOLD, bold=True, first=True, align=PP_ALIGN.CENTER)
    tf2 = textbox(s, Inches(1), Inches(3.3), Inches(11.333), Inches(2.0))
    first = True
    for ln in lines:
        add_para(tf2, ln, 16, WHITE if first else MUTED, first=first,
                 align=PP_ALIGN.CENTER, space_after=8)
        first = False
    return s


def scene_slide(kicker, title, img_name, groups, footer=None):
    """Big scene artwork with numbered position markers + keyed legend.
       groups: list of (heading, color, [(n, name, desc, fx, fy), ...]);
       fx/fy are 0..1 fractions of the rendered image rect."""
    s = prs.slides.add_slide(BLANK); add_bg(s)
    _kicker(s, kicker)
    tf = textbox(s, Inches(0.7), Inches(0.72), Inches(12), Inches(0.6))
    add_para(tf, title, 26, GOLD, bold=True, first=True)
    entry, _ = lookup(img_name)
    bx, by, bw, bh = Inches(0.6), Inches(1.55), Inches(7.9), Inches(5.27)
    rect(s, bx, by, bw, bh, CARDBG, rounded=True)
    img = fetch_img(entry) if entry else None
    rx0, ry0, rw0, rh0 = int(bx), int(by), int(bw), int(bh)
    if img:
        pad = int(Inches(0.1))
        ib_x, ib_y = int(bx) + pad, int(by) + pad
        ib_w, ib_h = int(bw) - 2 * pad, int(bh) - 2 * pad
        add_image_fit(s, img, ib_x, ib_y, ib_w, ib_h)
        with Image.open(img) as _im:
            pw, ph = _im.size
        scale = min(ib_w / pw, ib_h / ph)
        rw0 = int(pw * scale); rh0 = int(ph * scale)
        rx0 = ib_x + (ib_w - rw0) // 2; ry0 = ib_y + (ib_h - rh0) // 2
    D = int(Inches(0.34))
    for heading, color, rows in groups:
        for n, name, desc, fx, fy in rows:
            cx = rx0 + int(fx * rw0); cy = ry0 + int(fy * rh0)
            sp = s.shapes.add_shape(MSO_SHAPE.OVAL, cx - D // 2, cy - D // 2, D, D)
            sp.fill.solid(); sp.fill.fore_color.rgb = color
            sp.line.color.rgb = WHITE; sp.line.width = Pt(1.5)
            sp.shadow.inherit = False
            t2 = sp.text_frame
            t2.margin_top = 0; t2.margin_bottom = 0
            t2.margin_left = 0; t2.margin_right = 0; t2.word_wrap = False
            pp = t2.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
            rr = pp.add_run(); rr.text = str(n)
            rr.font.size = Pt(12); rr.font.bold = True; rr.font.color.rgb = WHITE
            _set_font(rr)
    lx = Inches(8.8); lw = Inches(4.4)
    tf = textbox(s, lx, Inches(1.5), lw, Inches(5.4))
    first = True
    for heading, color, rows in groups:
        add_para(tf, heading, 14, color, bold=True, first=first, space_after=4)
        first = False
        for n, name, desc, fx, fy in rows:
            add_runs(tf, [(str(n) + '  ', color, True), (name, WHITE, True),
                          ('  ' + desc, MUTED, False)], 12, space_after=3)
    if footer:
        tf = textbox(s, Inches(0.6), Inches(6.95), Inches(8.0), Inches(0.5))
        add_para(tf, footer, 10.5, MUTED, first=True, space_after=0)
    return s


# ================= BUILD DECK =================
title_slide(
    '大巴扎趣闻录',
    '物品背景、翻译彩蛋与文化梗（归属与效果均取自数据库）',
    '数据来源：bazaardb.gg · Season 16.1 · 卡面图片来自游戏数据库',
)

# ---- 海盗 Vanessa ----
section_slide('一', '海盗的航海世界', '瓦妮莎 · Vanessa')
item_slide('海狗沙龙', '海盗文化 · 航海俚语', [
    '“海狗”（Seadog）是航海俚语里饱经风霜的老水手；“沙龙”（Saloon）则是 18 世纪海盗聚众喝酒、分赃、交换情报的场所。',
    '机制“己方每有一件伙伴物品便 +1 多重释放”完美呼应——伙伴越多，沙龙越热闹，连锁反应越猛。',
    '加速一件、减速一件，恰如酒吧的嘈杂：有人越喝越亢奋，有人越喝越昏沉。',
])
item_slide('船首像', '海盗文化 · 航海传统', [
    '船首像是安在船头的守护雕像，古希腊人相信它能驱邪避灾、庇佑航程，也是一艘船的身份象征。',
    '效果呼应它在船头的位置：左侧（船前）水系物品冷却缩短，右侧（船后）物品获得属性加成——前后皆护。',
])
item_slide('潜水配重', '海盗文化 · 潜水装备', [
    '潜水配重是潜水员系在腰带上的铅块，用来抵消浮力顺利下潜，现实中按需分成多节增减。',
    '三层效果层层递进，最妙的是“多重释放增幅等于弹药量”——铅块（弹药）越多，配重越强，设计与实物严丝合缝。',
])
item_slide('旗舰', '海盗船员 · 经典座舰', [
    '旗舰（Flagship）是海盗职业的经典座舰，也是整支船组的核心——舰队司令坐镇的那艘船。',
    '效果“每拥有一种工具、地产、伙伴、弹药或遗物就 +1 多重释放”，正是这艘旗舰在号令满船的船员与物资：家底越厚，齐射越猛。',
])
grid_slide('海盗船员 · 旗舰船组', '不是“伙伴”，是你的船员', 'Friend 标签之下，其实是旗舰上的水手',
    ['理查森先生', '鹦鹉皮特', '橘利安', '盐钳海盗', '三花'],
    {'理查森先生': 'Mr. Richardson——称呼正式的老船员，触发加速/减速时自我强化。',
     '鹦鹉皮特': '海盗肩头鹦鹉的经典形象，Pesky=聒噪烦人；会飞起来灼烧敌人。',
     '橘利安': 'Orange Julian 谐音美国饮品连锁 Orange Julius，按累计金币为己方加伤。',
     '盐钳海盗': 'Old Salt=饱经风霜的老水手，配上钳爪如老海怪，也是能挥砍的武器。',
     '三花': '三花猫（Calico）是航海传统里的“船猫”——既能捕鼠又被水手视作幸运，招财猫亦多以三花为原型。'})
item_slide('龙涎香', '海盗文化 · 海上浮金', [
    '龙涎香是抹香鲸肠道的分泌物，排出后在海面漂浮、氧化陈化多年，才成为珍贵的定香香料。',
    '历史上它价比黄金，水手若在海上偶然拾得便可一夜暴富，故有“海上浮金”之称。',
])
item_slide('巨龟托图加', '海盗文化 · 海盗据点', [
    '托图加（Tortuga）是加勒比海伊斯帕尼奥拉岛北面的小岛，17 世纪臭名昭著的海盗据点与走私天堂。',
    '因岛形酷似一只伏卧的海龟，西班牙人便以“tortuga”（龟）为它命名。',
])
grid_slide('海盗文化 · 航海器物', '导航仪器与船上梗', '均为海盗物品，名字都有出处',
    ['六分仪', '星盘', '逞威风腰带扣'],
    {'六分仪': '大航海时代的定位神器，测量天体与地平线夹角推算纬度。',
     '星盘': '源自古希腊的天文测量仪，比六分仪更古老，是“星辰导航”的鼻祖。',
     '逞威风腰带扣': 'Swashbuckler=耀武扬威的剑客，词源 swash（挥剑声）+ buckler（小盾），拆成“腰带扣”双关。'})
item_slide('狼筅', '海盗武器 · 中华冷兵器', [
    '狼筅（Langxian）的英文名直接用了拼音，来源是中国明代的竹制长矛。',
    '它由一整根带枝桠的老竹削制、前端加装铁刃，是抗倭名将戚继光在“鸳鸯阵”中发扬光大的兵器：以密集枝桠扰乱敌阵、掩护刀手推进。',
    '有意思的是，这件当年专门对付倭寇（海盗）的武器，如今成了海盗自己的专属兵器；效果“每赢一场战斗便 +伤害”，也呼应它越战越强的实战传承。',
])
item_slide('淬锋钢', '海盗工具 · 修刃钢', [
    '淬锋钢（Honing Steel）就是厨房常见的“磨刀棒”，也叫修刃钢、厨师钢——一根用于矫正刀刃的金属棒。',
    '它并不磨掉刀上的金属，而是通过按压重塑，把日常使用中卷曲、倾斜的刃口重新扶正，瞬间恢复锐利。',
    '对应效果“己方最左侧和最右侧的武器获得属性加成”——它不造伤，而是替两端的武器“扶正刃口、磨砺锋芒”。',
])
item_slide('宠物石', '海盗玩具 · 流行文化梗', [
    '宠物石是 1975 年美国广告人加里·达尔的恶搞发明：一块普通鹅卵石，装进带透气孔的纸盒，还附赠一本“训练手册”当宠物卖。',
    '短短半年狂销约 150 万个，成为泡沫式流行文化的经典符号。',
    '游戏里“武器 + 伙伴 + 玩具”三合一标签，把这份荒诞感还原得淋漓尽致。',
])
item_slide('套娃', '海盗玩具 · 民俗器物', [
    '俄罗斯套娃（Matryoshka）诞生于 1890 年代，一个套一个、层层相扣，象征母性、生育与家族传承。',
    '游戏里它作为玩具物品出现，和现实一样——拆开一个，里面还有一个。',
])
item_slide('烈酒杯', '海盗器物 · 酒馆文化', [
    '烈酒杯是盛放一口烈酒（约 30–45 毫升）的小玻璃杯。“shot”的词源众说纷纭：一说源自旧西部以一发子弹（shot）换一杯威士忌，一说单指一“注”烈酒。',
    '对海盗与水手而言，朗姆酒配给是航海生活的日常，这只小杯自带浓浓海味。',
    '效果“加速 4 件、减速 4 件”恰似酒过三巡——同桌有人越喝越亢奋，有人越喝越犯困。',
])
item_slide('鱼雷', '海盗弹药 · 词源冷知识', [
    'Torpedo（鱼雷）一词借自会放电麻痹猎物的“电鳐”（torpedo 鱼）；现代自航鱼雷由罗伯特·怀特黑德于 1866 年发明。',
    '效果会随触发物品的大小翻倍加伤——这里“该物品”指的是触发它的那件水系/弹药物品，而非鱼雷自身。',
])

# 海盗的船与军械库
grid_slide('海盗的船 · 拆解一艘海盗船', '一艘海盗船由哪些部件组成', '船舵掌向、船锚泊定、鸦巢瞭望——船长舱作为核心舱室见下页详解',
    ['船舵', '船锚', '鸦巢'],
    {'船舵': '船上掌舵之处，加速相邻物品；若己方拥有载具或大型物品，其冷却减半。',
     '船锚': '起锚与抛锚，按敌人最大生命的百分比造成伤害。',
     '鸦巢': '桅杆顶端的瞭望台，水手在此放哨观敌；大幅提升己方武器暴击率。'})
item_slide('船长舱', '海盗的船 · 船长的居所', [
    '船长舱是船长起居与指挥的核心舱室，多设在船尾楼，通常是全船最宽敞、装潢最考究的房间。',
    '作为大型地产，它统筹全船：加速工具与载具、为己方装填弹药并整体强化属性，正如船长在此调度全员。',
    '下一页就走进这间船长舱，数一数画面里到底藏着多少真实存在的游戏物品。',
])
grid_slide('海盗的军械库 · 经典枪械', '海盗最标志性的火器', '从喇叭铳到燧发枪，海盗的“枪械谱”',
    ['雷筒', '火枪', '左轮手枪', '刺刀手枪', '吹箭枪', '大钢弩'],
    {'雷筒': '喇叭口散弹铳，17–18 世纪海盗登船近战的头号招牌枪。',
     '火枪': '前装滑膛枪，大航海时代陆海通用的主力火器。',
     '左轮手枪': '转轮连发手枪，象征海盗浪漫的经典短枪。',
     '刺刀手枪': '手枪与短刀合体，一枪不中立刻转白刃。',
     '吹箭枪': '无声吹管毒箭，源自原住民的狩猎技艺。',
     '大钢弩': '重型十字弩，严格说是弩而非火枪，胜在蓄力重击。'})
scene_slide('海盗的船 · 船长舱识图', '用编号标出画面里真实存在的物品', '船长舱', [
    ('航海与装备（真实海盗物品）', SEA, [
        (1, '侦察望远镜', '左上支架单筒镜', 0.13, 0.20),
        (2, '潜水头盔', '左侧桌上铜盔', 0.11, 0.45),
        (3, '六分仪／星盘', '书桌上的黄铜航海仪', 0.45, 0.51),
        (4, '通缉海报', '窗右墙上裱框肖像', 0.62, 0.36),
    ]),
    ('武器墙（真实海盗枪械）', GOLD, [
        (5, '火枪', '横挂长枪', 0.85, 0.22),
        (6, '左轮手枪', '枪架顶部转轮枪', 0.77, 0.27),
        (7, '大钢弩', '红色十字弩', 0.88, 0.36),
        (8, '刺刀手枪', '枪架中部燧发手枪', 0.80, 0.40),
        (9, '雷筒', '喇叭口短铳', 0.75, 0.49),
        (10, '吹箭枪', '弩下彩色细管', 0.91, 0.45),
    ]),
    ('家族象征', RGBColor(0xE0, 0x5A, 0x6A), [
        (11, '蔻森娜纹章', '窗上红白徽记·瓦妮莎家徽', 0.40, 0.26),
    ]),
], footer='注：画面里的地球仪属皮格马利安，镜子、宝箱、金币等仅为场景装饰；此前误标的“船首像”实为桌上的装饰镜——船首像只会立于船头，不会出现在船长舱内。')

# ---- 杜利 Dooley ----
section_slide('二', '杜利的翻译艺术', '杜利 · Dooley（机器人角色，命名充满谐音与双关）')
item_slide('墙莱士', '文字游戏 · 谐音双关', [
    '英文名 Wall-ez = Wall（墙）+ 人名后缀 -ez，直读近似“华莱士”（Wallace）。',
    '中文把“华”换成“墙”，既保住 Wall 的本义，又留住 Wallace 的发音——而它的属性正是防御型的护盾墙，名副其实。',
])
item_slide('内存狂暴模块', '文字游戏 · 双关', [
    'RAMpage 是绝妙双关：RAM 是计算机内存，Rampage 是暴走狂暴，合起来即“内存暴走”。',
    '中文“内存·狂暴模块”把两层含义都拆了出来；效果在触发减速、剧毒或冻结时爆发增益，正应“狂暴”之名。',
])
item_slide('加密货币', '文字游戏 · 赛博主题', [
    'Crypto 既指加密货币，也暗指密码学（cryptography），是机器人杜利科技谱系里的应景物品。',
    '把时下最热的数字货币做成卡牌，命名直白又贴合杜利的赛博主题。',
])
item_slide('木马病毒', '文字游戏 · 赛博主题', [
    'Trojan 一语双关：既是希腊神话里的特洛伊木马，也是计算机领域的“木马病毒”。',
    '以“伪装潜入、内部发作”的意象呼应杜利的机械与黑客风格。',
])
grid_slide('命名彩蛋 · 人名系列', '推土机化作一群工友', '英文名 = 现实事物 + 常见人名',
    ['推土比尔', '哈姆锤特', '岛弹小姐', '急救艾登', '微波戴夫', '铜蛇艾德', '饮水沃特', '空调皮埃尔'],
    {'推土比尔': 'Bill Dozer ← Bulldozer 推土机 + 人名 Bill。',
     '哈姆锤特': 'Hammlet ← Hammer 锤子 + 莎剧 Hamlet 哈姆雷特。',
     '岛弹小姐': 'Miss Isles ← Missile 导弹（谐音）+ Isles 岛屿。',
     '急救艾登': 'First Aiden ← First Aid 急救 + 人名 Aiden。',
     '微波戴夫': 'Micro Dave ← Microwave 微波炉 + 人名 Dave。',
     '铜蛇艾德': 'Copper Ed 谐音 Copperhead 铜头蛇。',
     '饮水沃特': 'Walter Cooler ← Water Cooler 饮水机 + Walter。',
     '空调皮埃尔': 'Pierre Conditioner ← Air Conditioner 空调 + Pierre。'})
grid_slide('命名彩蛋 · 恐龙家族', '每只恐龙都有名字', '英文名 = 名字 + Dinosaur 词根',
    ['妈妈暴龙', '小黛暴龙', '巨魔龙', '特里翼龙', '坦奇甲龙'],
    {'妈妈暴龙': 'Momma-Saur ← Momma 妈妈 + -saur 恐龙词根。',
     '小黛暴龙': 'Diana-Saur ← Diana 黛安娜 + Dinosaur。',
     '巨魔龙': 'Trollosaur ← Troll 巨魔 + -saur。',
     '特里翼龙': 'Terry-Dactyl ← Terry + Pterodactyl 翼龙。',
     '坦奇甲龙': 'Tanky Anky ← Tanky 坦克型 + Ankylosaurus 甲龙。'})
grid_slide('命名彩蛋 · 机械昆虫', '黑客语（Leet Speak）命名', '用数字替代字母：3=E、1=I、4=A…',
    ['机械红焰萤', '机械绿马蜂', '机械蓝甲虫', '机械黄螳螂', '机械黑蜘蛛'],
    {'机械红焰萤': 'RED-F1R3FLY = Red Firefly（红萤火虫）。',
     '机械绿马蜂': 'GRN-W4SP = Green Wasp（绿马蜂）。',
     '机械蓝甲虫': 'BLU-B33TL3 = Blue Beetle（蓝甲虫）。',
     '机械黄螳螂': 'YLW-M4NT1S = Yellow Mantis（黄螳螂）。',
     '机械黑蜘蛛': 'BLK-SP1D3R = Black Spider（黑蜘蛛）。'})
grid_slide('命名彩蛋 · 更多文字游戏', '科技与生活的混搭', '均为杜利物品',
    ['克里斯军刀', '砖友', '机能棒', '多尔王', '炫光 LED'],
    {'克里斯军刀': 'Chris Army Knife ← Swiss Army Knife 瑞士军刀，Swiss→Chris。',
     '砖友': 'Brick Buddy ← 砖头 Brick + 伙伴 Buddy，直白可爱。',
     '机能棒': 'Gearnola Bar ← Gear 齿轮 + Granola 麦片棒。',
     '多尔王': 'Dooltron ← Dooley + -tron 电子音效后缀，霸气外露。',
     '炫光 LED': 'Cool LEDs 字面“酷炫 LED”，既是外观又是科技梗。'})

# ---- 朱尔斯 Jules ----
section_slide('三', '朱尔斯的厨房', '朱尔斯 · Jules（多重释放与料理主题）')
item_slide('搅拌机', '厨房器械 · 料理链条', [
    'Blender（搅拌机）是朱尔斯厨房体系的核心器械，负责把食材打成饮品。',
    '它产出的正是“冰沙”（Slushee）——器械与产物名称一一对应，构成朱尔斯“食材—料理—产物”的完整链条。',
])
item_slide('冰沙', '厨房产物 · 消耗品', [
    'Slushee（冰沙）是搅拌机运转后产出的消耗品，直译“冰沙/沙冰”。',
    '在朱尔斯的厨房里，它是把原料转化为战力的“成品”，料理主题一以贯之。',
])

# ---- 皮格马利安 Pygmalien ----
section_slide('四', '皮格马利安的商业帝国', '皮格马利安 · Pygmalien（金钱、地产与冻结流）')
item_slide('耶塔里亚大棒', '皮格马利安 · 招牌武器', [
    '大型金质武器，触发冻结时充能。',
    '半血以下可将一件物品冻结 99 秒——在快节奏对局里近乎“永久”，是残局锁死对手关键物品的翻盘神器。',
])
item_slide('财源炮', '皮格马利安 · 金钱主题', [
    'Cash Cannon 直译“现金大炮”，把“钱”本身做成武器——用钞票开火。',
    '这份“财大气粗”的荒诞感，正是皮格马利安商业帝国主题的集中体现。',
])
grid_slide('皮格马利安 · 地产与广告', '越大越有排面', '均为皮格马利安物品',
    ['冰淇淋车', '广告牌', '摩天大楼'],
    {'冰淇淋车': 'Ice Cream Truck——街头循环叫卖的冰淇淋车，满满的市井商业气息。',
     '广告牌': 'Billboard——巨型户外广告牌，皮格马利安营销帝国的象征。',
     '摩天大楼': 'Skyscraper——直插云霄的摩天楼，“越大越有排面”的地产梗。'})

# ---- 斯特尔 Stelle ----
section_slide('五', '斯特尔的怪炮', '斯特尔 · Stelle（搞笑又硬核的火炮）')
item_slide('鸡农炮', '斯特尔 · 冷知识', [
    '别被名字骗了：“Chicken Cannon”（鸡炮）在现实中真实存在——航空业用它把鸡发射向机身，测试挡风玻璃与发动机的抗鸟撞能力。',
    '游戏里它造成等同相邻小/中型物品护盾之和的伤害，并让相邻物品“开始飞行”，堪称最“社死”的武器。',
])
item_slide('湮灭火炮', '斯特尔 · 硬控', [
    'Oblivion Cannon（湮灭火炮）直接摧毁敌人手中最小的物品，并造成 200 伤害。',
    '“直接删除”对手物品，是大巴扎里最不讲道理的控制手段之一。',
])

# ---- 马克 Mak ----
section_slide('六', '马克的炼金术', '马克 · Mak（药水、炼金与冷门译名）')
item_slide('炼金梨缶', '马克 · 炼金古器', [
    'Aludel（阿吕德尔）是中世纪炼金术使用的梨形陶土容器，一节节叠起，用于加热升华物质并收集升华产物。',
    '中文“梨缶”既描摹其梨形外观，又保留了古器皿的韵味，是相当讲究的冷门译名。',
])
grid_slide('马克 · 更多炼金器物', '冷门却讲究的译名', '均为马克物品',
    ['瘟疫长柄刀', '驼鹿角杖'],
    {'瘟疫长柄刀': 'Plague Glaive——glaive 是欧洲的长柄战刀，并非 Halberd（戟）。',
     '驼鹿角杖': 'Staff of the Moose——顶端装饰驼鹿角的法杖，造型粗犷。'})

closing_slide('趣闻待续', [
    '归属与效果说明均取自数据库（Season 16.1），卡面图片来自游戏数据库。',
    '数据来源：bazaardb.gg · 欢迎补充更多物品彩蛋',
])

# ================= SAVE =================
out = os.path.join(BASE, 'output', '大巴扎趣闻录.pptx')
prs.save(out)
print('SAVED', out, 'slides=', len(prs.slides._sldIdLst))
